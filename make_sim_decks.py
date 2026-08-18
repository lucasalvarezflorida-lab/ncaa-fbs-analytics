"""Generate the Season Sim PowerPoint set:
  decks/2026_Season_Sim_Overview.pptx        (national deck — safe to commit)
  decks/2026_<Conference>_Projections.pptx   (one per conference, x10:
      standings + title race + TWO SLIDES PER TEAM — projections/strengths/
      weaknesses profile, and the projected depth chart)

LOCAL-ONLY WARNING: the conference decks embed OurLads depth charts
(ourlads_depth.json is gitignored — their curated product). The conference
decks are gitignored for the same reason; only the overview deck is tracked.

Data = the same engines as the workbook: independent per-team sim for the
national projections (Season Sim tab numbers) and the joint
_simulate_conf_standings() for conference standings/title races.
Brand: workbook navy 0A2851 / orange F47321, Arial throughout.
"""

import json
import sys
from pathlib import Path

import numpy as np

MODEL = Path(r"C:\Users\lucas\Fun Projects\Sports Data Analysis\ncaa-fbs-model")
sys.path.insert(0, str(MODEL))
sys.path.insert(0, str(MODEL / "fpi-decomposition"))

import cfbd_client as cfbd
from build_conference_book import (CONF_ORDER, HFA, SIM_SIGMA, UNRATED_MARGIN,
                                   _margin_model, _simulate_conf_standings,
                                   load_fpi_2026, load_scouting)
from name_mapping import normalize_name as norm

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0A, 0x28, 0x51)
ORANGE = RGBColor(0xF4, 0x73, 0x21)
ICE = RGBColor(0xEA, 0xF0, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x27, 0x3D)
MUTE = RGBColor(0x5C, 0x6B, 0x7E)
PEACH = RGBColor(0xFD, 0xE9, 0xD9)

DECKS = MODEL / "decks"
DECKS.mkdir(exist_ok=True)
N_SIMS = 10000

# ---------------- data ----------------
_fpi_raw = load_fpi_2026()
fpi = {k: (v if isinstance(v, dict) else {"fpi": float(v)})
       for k, v in _fpi_raw.items()}


def rate(team):
    e = fpi.get(norm(team))
    return float(e["fpi"]) if e else None


def pick(g, *names):
    for nm in names:
        if g.get(nm) is not None:
            return g[nm]
    return None


raw = cfbd.get("/games", {"year": 2026, "seasonType": "regular"})
games, team_conf = [], {}
for g in raw:
    h, a = pick(g, "homeTeam", "home_team"), pick(g, "awayTeam", "away_team")
    hc = pick(g, "homeConference", "home_conference")
    ac = pick(g, "awayConference", "away_conference")
    hk = str(pick(g, "homeClassification", "home_classification")).lower()
    ak = str(pick(g, "awayClassification", "away_classification")).lower()
    CONF_ALIAS = {"American Athletic": "American"}
    if hk == "fbs" and hc:
        team_conf[h] = CONF_ALIAS.get(hc, hc)
    if ak == "fbs" and ac:
        team_conf[a] = CONF_ALIAS.get(ac, ac)
    games.append(dict(home=h, away=a,
                      neutral=bool(pick(g, "neutralSite", "neutral_site")),
                      conf_game=bool(pick(g, "conferenceGame", "conference_game"))))

# independent per-team sim (Season Sim tab numbers)
win_prob, resid, margin_desc = _margin_model()
SIM_SD = f"{float(resid.std()):.1f}" if resid is not None else f"{SIM_SIGMA:g}"
rng = np.random.default_rng(2026)
win_probs = {t: [] for t in team_conf}
for g in games:
    for team, opp, is_home in ((g["home"], g["away"], True),
                               (g["away"], g["home"], False)):
        if team not in win_probs:
            continue
        rt, ro = rate(team), rate(opp)
        if rt is None:
            continue
        margin = UNRATED_MARGIN if ro is None else rt - ro
        if not g["neutral"]:
            margin += HFA if is_home else -HFA
        win_probs[team].append(win_prob(margin))

national = []
for t, ps in win_probs.items():
    if not ps:
        continue
    sims = (rng.random((N_SIMS, len(ps))) < np.array(ps)).sum(axis=1)
    national.append(dict(team=t, conf=team_conf[t], fpi=rate(t),
                         games=len(ps), mean=float(sims.mean()),
                         ten=float((sims >= 10).mean()),
                         bowl=float((sims >= 6).mean())))
national.sort(key=lambda d: -d["mean"])

conf_data = {}
for conf in CONF_ORDER:
    rows, split = _simulate_conf_standings(
        games, fpi, team_conf, conf, N_SIMS, np.random.default_rng(2026),
        resid=resid)
    if rows is not None:
        conf_data[conf] = (rows, split)

national_map = {d["team"]: d for d in national}
scouting = load_scouting().get("teams", {})
_depth_raw = json.loads((MODEL / "ourlads_depth.json").read_text(
    encoding="utf-8"))
depth_by_norm = {norm(k): v for k, v in _depth_raw.get("teams", {}).items()}
DEPTH_STAMP = _depth_raw.get("captured", "")
print(f"data ready: {len(national)} teams, {len(conf_data)} conferences, "
      f"depth for {len(depth_by_norm)} teams ({DEPTH_STAMP})")

# OurLads rows arrive ordered offense -> defense -> special teams; split on
# the first defensive / first ST position.
DEF_POS = {"DE", "DT", "NT", "NG", "DL", "EDGE", "RUSH", "JACK", "BANDIT",
           "LB", "ILB", "OLB", "MLB", "MIKE", "WILL", "SAM", "CB", "S", "FS",
           "SS", "DB", "NICKEL", "NB", "N", "STAR", "ROVER", "HUSKY", "SPUR",
           "N/S", "LEO", "STUD"}
ST_POS = {"K", "P", "PK", "LS", "KR", "PR", "KO", "K/P", "PT"}


def _is_def_pos(base, full):
    if full in DEF_POS or base in DEF_POS:
        return True
    # side-prefixed variants: LDE/RDE, LCB/RCB, LDT/RDT, LOLB/ROLB ...
    return len(base) > 2 and base[0] in "LR" and base[1:] in DEF_POS


def split_depth(rows):
    off, dfn, st = [], [], []
    seen_def = False
    for r in rows:
        base = str(r.get("pos", "")).split("-")[0].split("/")[0].upper()
        full = str(r.get("pos", "")).upper()
        if full in ST_POS or base in ST_POS:
            st.append(r)
        elif base == "H" and seen_def:
            continue  # holder — duplicate of the punter, skip
        elif _is_def_pos(base, full):
            seen_def = True
            dfn.append(r)
        elif seen_def:
            dfn.append(r)  # unknown position after the defense started
        else:
            off.append(r)
    return off, dfn, st


def depth_table(slide, title, rows, x, y, w, font=9.5):
    txt(slide, x, y, w, 0.3, title, 13, NAVY, bold=True)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, Inches(x), Inches(y + 0.35),
                                   Inches(w), Inches(0.26 * nrows))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, (h, cw) in enumerate(zip(["Pos", "Starter", "Second"],
                                    [0.75, 2.6, 2.6])):
        tbl.columns[j].width = Emu(int(Inches(cw * w / 5.95)))
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j else PP_ALIGN.CENTER
        f = p.runs[0].font
        f.name = "Arial"
        f.size = Pt(font)
        f.bold = True
        f.color.rgb = WHITE
    for i, r in enumerate(rows, 1):
        players = r.get("players") or []
        vals = [str(r.get("pos", "")),
                players[0] if players else "",
                players[1] if len(players) > 1 else ""]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = ICE if i % 2 == 0 else WHITE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j else PP_ALIGN.CENTER
            if p.runs:  # empty cells (no backup listed) have no run
                f = p.runs[0].font
                f.name = "Arial"
                f.size = Pt(font)
                f.color.rgb = INK
    for r_ in tbl.rows:
        r_.height = Emu(int(Inches(0.24)))
    return shape


def _fit_size(text, width_in, height_in, sizes=(13, 12, 11, 10, 9.5, 9)):
    """Largest Arial size whose wrapped line count fits the box."""
    paras = [p for p in str(text).split("\n")]
    for s in sizes:
        chars_per_line = max(20, int(width_in * 72 / (s * 0.52)))
        lines = sum(max(1, -(-len(p) // chars_per_line)) for p in paras)
        if lines * s * 1.25 <= height_in * 72:
            return s
    return sizes[-1]


def deep_block(slide, x, y, w, h, label, text):
    txt(slide, x, y, w, 0.35, label, 14, NAVY, bold=True)
    body = str(text or "").strip()
    if not body:
        body = "(none)"
    if len(body) < 60:
        # terse one-liner (G6 predictions): render as a callout
        txt(slide, x, y + 0.55, w, 1.2, body, 30, ORANGE, bold=True)
        return
    size = _fit_size(body, w, h - 0.55)
    txt(slide, x, y + 0.5, w, h - 0.5, body, size)


def team_slides(prs, conf, d):
    """Four slides per team: profile, deep dive x2, depth chart."""
    team = d["team"]
    sc = scouting.get(team) or {}
    nat = national_map.get(team) or {}

    # ---- slide A: projections + strengths/weaknesses ----
    s = blank(prs)
    txt(s, 0.9, 0.4, 9.0, 0.7, team, 32, NAVY, bold=True)
    txt(s, 0.9, 1.05, 9.0, 0.4, f"{conf} · ESPN 2026 preseason FPI "
        f"{d['fpi']:.1f}", 13, MUTE)
    cards = [(f"{d['over_w']:.1f}–{d['over_l']:.1f}", "proj overall"),
             (f"{d['conf_w']:.1f}–{d['conf_l']:.1f}", "proj conference"),
             (f"{nat.get('bowl', 0):.0%}", "P(bowl, 6+ wins)"),
             (f"{d['cg']:.0%}", "P(title game)"),
             (f"{d['champ']:.0%}", "P(champion)")]
    for i, (big, label) in enumerate(cards):
        stat_card(s, 0.7 + i * 2.45, 1.8, 2.3, big, label)
    sw_y = 3.6
    txt(s, 0.9, sw_y, 5.6, 0.4, "STRENGTHS", 14, NAVY, bold=True)
    str_list = sc.get("s") or ["(no scouting card)"]
    txt(s, 0.9, sw_y + 0.45, 5.6, 2.2,
        "\n".join(f"+  {x}" for x in str_list), 13)
    txt(s, 6.9, sw_y, 5.6, 0.4, "WEAKNESSES", 14, NAVY, bold=True)
    wk_list = sc.get("w") or ["(no scouting card)"]
    txt(s, 6.9, sw_y + 0.45, 5.6, 2.2,
        "\n".join(f"–  {x}" for x in wk_list), 13)
    scheme_bits = []
    if sc.get("ob"):
        scheme_bits.append(f"Offense: {sc['ob']}")
    if sc.get("db"):
        scheme_bits.append(f"Defense: {sc['db']}")
    if scheme_bits:
        txt(s, 0.9, 6.35, 11.6, 1.0, "\n".join(scheme_bits), 11, MUTE,
            italic=True)

    # ---- slides B/C: the deep dive ----
    deep = sc.get("deep") or {}
    s = blank(prs)
    txt(s, 0.9, 0.35, 11.6, 0.6, f"{team} — deep dive", 28, NAVY, bold=True)
    txt(s, 0.9, 0.98, 11.6, 0.35, "The season · prediction", 12, MUTE,
        italic=True)
    deep_block(s, 0.9, 1.55, 5.7, 5.5, "THE SEASON", deep.get("intro"))
    deep_block(s, 6.9, 1.55, 5.7, 5.5, "PREDICTION", deep.get("prediction"))

    s = blank(prs)
    txt(s, 0.9, 0.35, 11.6, 0.6, f"{team} — deep dive", 28, NAVY, bold=True)
    txt(s, 0.9, 0.98, 11.6, 0.35, "Roster · identity", 12, MUTE, italic=True)
    deep_block(s, 0.9, 1.55, 5.7, 5.3, "ROSTER", deep.get("roster"))
    deep_block(s, 6.9, 1.55, 5.7, 5.3,
               "IDENTITY · STRENGTHS · WEAKNESSES · GOALS",
               deep.get("identity"))
    if deep.get("vintage"):
        txt(s, 0.9, 7.05, 11.6, 0.35, f"Prep vintage: {deep['vintage']}", 8.5,
            MUTE, italic=True)

    # ---- slide D: projected depth chart ----
    s = blank(prs)
    entry = depth_by_norm.get(norm(team))
    txt(s, 0.9, 0.35, 11.6, 0.6, f"{team} — projected depth chart", 28,
        NAVY, bold=True)
    if not entry:
        txt(s, 0.9, 1.6, 11.0, 1.0, "No depth chart captured.", 14, MUTE)
        return
    sub = (f"OurLads, updated {entry.get('updated', DEPTH_STAMP)} · "
           f"{entry.get('off_scheme', '')} / {entry.get('def_scheme', '')} · "
           f"local use only — do not redistribute")
    txt(s, 0.9, 0.95, 11.6, 0.35, sub, 11, MUTE, italic=True)
    off, dfn, st = split_depth(entry.get("rows") or [])
    font = 9.5 if max(len(off), len(dfn) + len(st) + 1) <= 14 else 8.5
    depth_table(s, "OFFENSE", off[:14], 0.7, 1.45, 5.95, font=font)
    depth_table(s, "DEFENSE", dfn[:14], 6.85, 1.45, 5.95, font=font)
    if st:
        y_st = 1.45 + 0.35 + 0.26 * (min(len(dfn), 14) + 1) + 0.25
        if y_st < 6.4:
            depth_table(s, "SPECIAL TEAMS", st[:4], 6.85, y_st, 5.95,
                        font=font)


print(f"team content wired: scouting {len(scouting)}, "
      f"depth {len(depth_by_norm)}")

# ---------------- pptx helpers ----------------


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def txt(slide, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
        italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = "Arial"
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return box


def stat_card(slide, x, y, w, big, label, big_color=ORANGE):
    txt(slide, x, y, w, 1.0, big, 40, big_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x, y + 0.75, w, 0.6, label, 12, MUTE, align=PP_ALIGN.CENTER)


def title_slide(prs, title, subtitle, kicker):
    s = blank(prs, NAVY)
    txt(s, 0.9, 2.2, 11.5, 0.5, kicker.upper(), 14, ORANGE, bold=True)
    txt(s, 0.9, 2.7, 11.5, 1.8, title, 44, WHITE, bold=True)
    txt(s, 0.9, 4.6, 11.5, 1.2, subtitle, 15, RGBColor(0xCA, 0xDC, 0xFC))
    return s


def standings_table(slide, rows, x, y, w, font=10.5):
    headers = ["Pos", "Team", "FPI", "Conf W–L", "Overall", "P(CG)", "P(Champ)"]
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, len(headers), Inches(x), Inches(y),
                                   Inches(w), Inches(0.32 * nrows))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    widths = [0.7, 3.1, 0.9, 1.35, 1.35, 1.0, 1.2]
    scale = w / sum(widths)
    for j, cw in enumerate(widths):
        tbl.columns[j].width = Emu(int(Inches(cw * scale)))
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
        f = p.runs[0].font
        f.name = "Arial"
        f.size = Pt(font)
        f.bold = True
        f.color.rgb = WHITE
    for i, d in enumerate(rows, 1):
        name = d["team"] + (" *" if d["ineligible"] else "")
        vals = [str(i), name, f"{d['fpi']:.1f}",
                f"{d['conf_w']:.1f}–{d['conf_l']:.1f}",
                f"{d['over_w']:.1f}–{d['over_l']:.1f}",
                f"{d['cg']:.0%}", f"{d['champ']:.0%}"]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = PEACH if i <= 2 else (
                ICE if i % 2 == 0 else WHITE)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
            f = p.runs[0].font
            f.name = "Arial"
            f.size = Pt(font)
            f.color.rgb = INK
            f.bold = (i <= 2 and j == 1)
    for r in tbl.rows:
        r.height = Emu(int(Inches(0.30)))
    return shape


def champ_chart(slide, rows, x, y, w, h, top=10):
    data = CategoryChartData()
    top_rows = sorted(rows, key=lambda d: -d["champ"])[:top]
    top_rows = [d for d in top_rows if d["champ"] >= 0.005]
    data.categories = [d["team"] for d in reversed(top_rows)]
    data.add_series("P(champ)", [d["champ"] for d in reversed(top_rows)])
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x),
                                    Inches(y), Inches(w), Inches(h), data)
    ch = gframe.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ORANGE
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = "0%"
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10)
    dl.font.name = "Arial"
    dl.font.color.rgb = INK
    cat_ax = ch.category_axis
    cat_ax.tick_labels.font.size = Pt(10)
    cat_ax.tick_labels.font.name = "Arial"
    val_ax = ch.value_axis
    val_ax.has_major_gridlines = False
    val_ax.visible = False
    return gframe


# ---------------- overview deck ----------------
prs = new_deck()
title_slide(
    prs, "2026 Season Simulation",
    "10,000 Monte Carlo seasons · joint conference standings · title races\n"
    "Prior: ESPN 2026 preseason FPI (July 14) · margins = FPI gap + 2.5 HFA "
    f"+ empirical residuals (σ {SIM_SD})\n"
    "Predictions frozen at kickoff of the 2026 opener — Aug 29",
    "NCAA FBS Analytics · projected standings")

# methodology + stat callouts
s = blank(prs)
txt(s, 0.9, 0.55, 11.5, 0.7, "How the simulation works", 36, NAVY, bold=True)
txt(s, 0.9, 1.6, 6.4, 4.6,
    "Every conference game is sampled once per season and shared between both "
    "teams, so the standings hang together — one team's upset is another's "
    "loss column.\n\n"
    "Standings run on conference win percentage (slates are 8 or 9 games). "
    "Two-way ties for the 1-seed break on the simulated head-to-head result; "
    "the title game is played at a neutral site.\n\n"
    "Transitioning programs (North Dakota State, Sacramento State) play their "
    "full league schedule but cannot take a championship-game seed.\n\n"
    "\u2018Split\u2019 = share of seasons in which the conference champion is "
    "NOT the team with the league's best overall record.", 15)
stat_card(s, 8.0, 1.7, 2.4, "10,000", "seasons per league")
stat_card(s, 10.5, 1.7, 2.4, f"σ {SIM_SD}", "empirical margin sd, 2021–25 fit")
stat_card(s, 8.0, 3.5, 2.4, "+2.5", "home-field points")
stat_card(s, 10.5, 3.5, 2.4, "138", "FBS teams rated")
txt(s, 8.0, 5.3, 4.9, 1.4,
    "2025 receipt: Duke won the ACC at 8-5 while 10-2 Miami stayed home — "
    "the champion \u2260 best team split is close to a coin flip in the "
    "biggest leagues.", 12, MUTE, italic=True)

# national title-race overview table
s = blank(prs)
txt(s, 0.9, 0.55, 11.5, 0.7, "The ten title races at a glance", 36, NAVY, bold=True)
headers = ["Conference", "Favorite", "P(champ)", "Projected CG pair", "Split"]
tbl_shape = s.shapes.add_table(len(conf_data) + 1, 5, Inches(0.9), Inches(1.5),
                               Inches(11.5), Inches(0.34 * (len(conf_data) + 1)))
tbl = tbl_shape.table
tbl.first_row = False
tbl.horz_banding = False
for j, (h, cw) in enumerate(zip(headers, [2.6, 2.6, 1.3, 3.6, 1.4])):
    tbl.columns[j].width = Emu(int(Inches(cw)))
    c = tbl.cell(0, j)
    c.text = h
    c.fill.solid()
    c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if j in (0, 1, 3) else PP_ALIGN.CENTER
    f = p.runs[0].font
    f.name = "Arial"
    f.size = Pt(12)
    f.bold = True
    f.color.rgb = WHITE
for i, conf in enumerate([c for c in CONF_ORDER if c in conf_data], 1):
    rows, split = conf_data[conf]
    fav = max(rows, key=lambda d: d["champ"])
    pair = " vs ".join(d["team"] for d in rows[:2])
    vals = [conf, fav["team"], f"{fav['champ']:.0%}", pair, f"{split:.0%}"]
    for j, v in enumerate(vals):
        c = tbl.cell(i, j)
        c.text = v
        c.fill.solid()
        c.fill.fore_color.rgb = ICE if i % 2 == 0 else WHITE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j in (0, 1, 3) else PP_ALIGN.CENTER
        f = p.runs[0].font
        f.name = "Arial"
        f.size = Pt(12)
        f.color.rgb = INK
        f.bold = (j == 1)

# projected top 25
s = blank(prs)
txt(s, 0.9, 0.55, 11.5, 0.7, "Projected Top 25 — overall wins", 36, NAVY,
    bold=True)
half = 13
for col, chunk in ((0, national[:half]), (1, national[half:25])):
    x = 0.9 + col * 6.0
    headers = ["#", "Team", "Conf", "FPI", "Proj W", "P(10+)"]
    tshape = s.shapes.add_table(len(chunk) + 1, 6, Inches(x), Inches(1.5),
                                Inches(5.7), Inches(0.31 * (len(chunk) + 1)))
    t = tshape.table
    t.first_row = False
    t.horz_banding = False
    for j, (h, cw) in enumerate(zip(headers, [0.5, 2.0, 1.1, 0.7, 0.8, 0.8])):
        t.columns[j].width = Emu(int(Inches(cw)))
        c = t.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j in (1, 2) else PP_ALIGN.CENTER
        f = p.runs[0].font
        f.name = "Arial"
        f.size = Pt(10)
        f.bold = True
        f.color.rgb = WHITE
    for i, d in enumerate(chunk, 1):
        rank = i + col * half
        conf_short = {"Conference USA": "CUSA", "Mid-American": "MAC",
                      "Mountain West": "MWC", "FBS Independents": "Ind",
                      "American Athletic": "American"}.get(d["conf"], d["conf"])
        vals = [str(rank), d["team"], conf_short, f"{d['fpi']:.0f}",
                f"{d['mean']:.1f}", f"{d['ten']:.0%}"]
        for j, v in enumerate(vals):
            c = t.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = ICE if i % 2 == 0 else WHITE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j in (1, 2) else PP_ALIGN.CENTER
            f = p.runs[0].font
            f.name = "Arial"
            f.size = Pt(10)
            f.color.rgb = INK
    for r in t.rows:
        r.height = Emu(int(Inches(0.29)))

# split chart slide
s = blank(prs)
txt(s, 0.9, 0.55, 11.5, 0.7, "The trophy \u2260 the best team", 36, NAVY,
    bold=True)
txt(s, 0.9, 1.35, 11.5, 0.5,
    "Share of simulated seasons where the champion does NOT have the league's "
    "best overall record", 14, MUTE)
data = CategoryChartData()
confs_sorted = sorted(conf_data.items(), key=lambda kv: kv[1][1])
data.categories = [c for c, _ in confs_sorted]
data.add_series("split", [v[1] for _, v in confs_sorted])
gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.2),
                            Inches(2.0), Inches(8.2), Inches(5.0), data)
ch = gframe.chart
ch.has_legend = False
ch.has_title = False
plot = ch.plots[0]
plot.gap_width = 60
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = NAVY
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format = "0%"
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(11)
dl.font.name = "Arial"
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = "Arial"
ch.value_axis.visible = False
ch.value_axis.has_major_gridlines = False
txt(s, 9.7, 2.6, 2.9, 3.6,
    "Multi-contender leagues split most: the SEC and Big Ten crown the "
    "\u2018wrong\u2019 team over half the time.\n\nOne-horse leagues (Pac-12: "
    "Boise) split least — the favorite usually holds both crowns.", 12, MUTE)

out = DECKS / "2026_Season_Sim_Overview.pptx"
prs.save(out)
print("wrote", out.name)

# ---------------- per-conference decks ----------------
for conf, (rows, split) in conf_data.items():
    prs = new_deck()
    fav = max(rows, key=lambda d: d["champ"])
    title_slide(
        prs, f"{conf} — 2026 Projections",
        f"Projected standings, title race, and a four-slide breakdown of all "
        f"{len(rows)} teams:\nprojections · strengths & weaknesses · full "
        f"deep dives · projected depth charts\n"
        f"10,000 joint simulated seasons · ESPN 2026 preseason FPI prior · "
        f"LOCAL USE (OurLads depth data)",
        "NCAA FBS Analytics · conference breakdown")

    # standings slide
    s = blank(prs)
    txt(s, 0.9, 0.4, 11.5, 0.7, f"{conf}: projected standings", 32, NAVY,
        bold=True)
    font = 10.5 if len(rows) <= 14 else 9
    standings_table(s, rows, 0.9, 1.35, 8.6, font=font)
    txt(s, 9.8, 1.5, 2.9, 4.8,
        f"Favorite: {fav['team']} ({fav['champ']:.0%})\n\n"
        f"Projected title game:\n{rows[0]['team']} vs {rows[1]['team']}\n\n"
        f"Split: champion \u2260 best overall record in {split:.0%} of "
        f"seasons.\n\nShaded rows = projected CG pair."
        + ("\n\n* title-ineligible (transition)" if any(d["ineligible"]
                                                        for d in rows) else ""),
        12, MUTE)

    # title race chart slide
    s = blank(prs)
    txt(s, 0.9, 0.4, 11.5, 0.7, f"{conf}: the title race", 32, NAVY, bold=True)
    champ_chart(s, rows, 1.0, 1.4, 8.0, 5.5, top=10)
    contenders = sorted(rows, key=lambda d: -d["champ"])
    c2 = contenders[1] if len(contenders) > 1 else None
    gap_team = max(rows, key=lambda d: (d["over_w"] / max(1, d["over_w"] + d["over_l"]))
                   - (d["conf_w"] / max(1, d["conf_w"] + d["conf_l"])))
    note = (f"{fav['team']} wins the league in {fav['champ']:.0%} of seasons"
            + (f"; {c2['team']} is next at {c2['champ']:.0%}." if c2 else "."))
    note += (f"\n\nSchedule note: {gap_team['team']}'s overall projection "
             f"({gap_team['over_w']:.1f} wins) outruns its conference finish — "
             f"draw matters.")
    txt(s, 9.4, 1.8, 3.3, 4.6, note, 12, MUTE)

    # team breakdowns: two slides per team, in projected-standings order
    for d in rows:
        team_slides(prs, conf, d)

    safe = conf.replace(" ", "_").replace("-", "_")
    out = DECKS / f"2026_{safe}_Projections.pptx"
    prs.save(out)
    print("wrote", out.name)

print("done:", len(conf_data) + 1, "decks in", DECKS)
