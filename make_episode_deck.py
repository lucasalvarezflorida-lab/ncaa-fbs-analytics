"""Episode 2 Week 1 deck — recap edition.
Slide 2 is the Week 0 receipts board (frozen v12 predictions graded vs
finals and the last pre-kick ledger pull). Per-game flow unchanged from
v12: what decides it -> keys per team -> score prediction last; flags and
lean lines retired, line movement promoted. Ep1's GAMES block is archived
below as _GAMES_EP1 (the frozen predictions live in git + the ledger).
Motif: real school logos (ESPN 500px PNGs in decks/logos/, URLs cached in
rosters/data/teams_fbs_2026.json) on white pucks over navy. Navy score-bug
panel per game with the model/market numbers and a win-probability split
bar. Dark title + dark 'The Card' closer around light content slides."""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0A, 0x28, 0x51)
NAVY2 = RGBColor(0x12, 0x35, 0x68)
ORANGE = RGBColor(0xF4, 0x73, 0x21)
ICE = RGBColor(0xEA, 0xF0, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x27, 0x3D)
MUTE = RGBColor(0x5C, 0x6B, 0x7E)
LIGHTLINE = RGBColor(0xD5, 0xDF, 0xEC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank(bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def txt(slide, x, y, w, h, text, size, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, anchor=None, spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = "Arial"
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return box


def shape(slide, kind, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def luminance(c):
    return 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]


LOGO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "decks", "logos")


def logo_badge(slide, x, y, d, key, plate=False):
    """School logo at (x, y), d inches square. plate=True adds a white
    circular puck behind it — needed on navy (UVA's logo is navy)."""
    if plate:
        shape(slide, MSO_SHAPE.OVAL, x, y, d, d, WHITE)
        inset = d * 0.10
    else:
        inset = 0.0
    slide.shapes.add_picture(
        os.path.join(LOGO_DIR, TEAMS[key]["logo"]),
        Inches(x + inset), Inches(y + inset),
        Inches(d - 2 * inset), Inches(d - 2 * inset))


def wp_bar(slide, x, y, w, h, team_a, pct_a, col_a, team_b, pct_b, col_b,
           label_color=WHITE):
    """Broadcast-style split probability bar."""
    wa = w * pct_a / 100.0
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, RGBColor(*col_b))
    sp = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, max(wa, 0.6), h,
               RGBColor(*col_a))
    txt(slide, x + 0.08, y + h / 2 - 0.11, 2.4, 0.25,
        f"{team_a} {pct_a:.0f}%", 10.5, WHITE if luminance(col_a) < 150 else INK,
        bold=True)
    txt(slide, x + w - 2.48, y + h / 2 - 0.11, 2.4, 0.25,
        f"{pct_b:.0f}% {team_b}", 10.5,
        WHITE if luminance(col_b) < 150 else INK, bold=True,
        align=PP_ALIGN.RIGHT)


TEAMS = {
    "UNC": dict(code="UNC", color=(0x7B, 0xAF, 0xD4), logo="unc.png"),
    "TCU": dict(code="TCU", color=(0x4D, 0x19, 0x79), logo="tcu.png"),
    "NCSU": dict(code="NCSU", color=(0xCC, 0x00, 0x00), logo="ncsu.png"),
    "UVA": dict(code="UVA", color=(0x23, 0x2D, 0x4B), logo="uva.png"),
    "JSU": dict(code="JSU", color=(0xB5, 0x12, 0x1B), logo="jsu.png"),
    "NDSU": dict(code="NDSU", color=(0x0A, 0x56, 0x40), logo="ndsu.png"),
    "HAW": dict(code="HAW", color=(0x00, 0x34, 0x20), logo="hawaii.png"),
    "STAN": dict(code="STAN", color=(0x8C, 0x15, 0x15), logo="stanford.png"),
    "MEM": dict(code="MEM", color=(0x00, 0x49, 0x91), logo="memphis.png"),
    "UNLV": dict(code="UNLV", color=(0xB1, 0x02, 0x02), logo="unlv.png"),
    "BAY": dict(code="BAY", color=(0x15, 0x45, 0x35), logo="baylor.png"),
    "AUB": dict(code="AUB", color=(0x0C, 0x24, 0x40), logo="auburn.png"),
    "CLEM": dict(code="CLEM", color=(0xF6, 0x67, 0x33), logo="clemson.png"),
    "LSU": dict(code="LSU", color=(0x46, 0x1D, 0x7C), logo="lsu.png"),
    "LOU": dict(code="LOU", color=(0xAD, 0x00, 0x00), logo="louisville.png"),
    "MISS": dict(code="MISS", color=(0x14, 0x21, 0x3D), logo="olemiss.png"),
    "WIS": dict(code="WIS", color=(0xC5, 0x05, 0x0C), logo="wisconsin.png"),
    "ND": dict(code="ND", color=(0x0C, 0x23, 0x40), logo="notredame.png"),
    "SMU": dict(code="SMU", color=(0x00, 0x33, 0xA0), logo="smu.png"),
    "FSU": dict(code="FSU", color=(0x78, 0x2F, 0x40), logo="fsu.png"),
}

NAME2CODE = {"North Carolina": "UNC", "TCU": "TCU", "NC State": "NCSU",
             "Virginia": "UVA", "Jacksonville State": "JSU",
             "North Dakota State": "NDSU", "Hawai'i": "HAW",
             "Stanford": "STAN", "Memphis": "MEM", "UNLV": "UNLV",
             "Baylor": "BAY", "Auburn": "AUB", "Clemson": "CLEM",
             "LSU": "LSU", "Louisville": "LOU", "Ole Miss": "MISS",
             "Wisconsin": "WIS", "Notre Dame": "ND", "SMU": "SMU",
             "Florida State": "FSU"}

# ---- card_data contract (review item A): market + model numbers come from
# edge_report.py --publish, never from hand-typed literals. Narrative fields
# (decides / honesty / keys) stay authored here; lean + players are kept as
# data but no longer rendered.
CARD_DATA = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                         "card_data_week1.json")


def load_card_data():
    if not os.path.exists(CARD_DATA):
        return {}, None
    import json
    d = json.load(open(CARD_DATA, encoding="utf-8"))
    return {(g["away"], g["home"]): g for g in d["games"]}, d.get("lines_as_of")


def _dash(x):
    """-7.5 -> '–7.5', 3 -> '+3' (home-perspective spread for display)."""
    return f"–{abs(x):g}" if x < 0 else (f"+{x:g}" if x > 0 else "PK")


def _fmt_ts(ts):
    import datetime as dt
    if not ts:
        return "?"
    t = dt.datetime.fromisoformat(ts.replace("Z", "+00:00"))
    t = t.astimezone(dt.timezone(dt.timedelta(hours=-4)))  # ET in August
    return f"{t:%b %d}".replace(" 0", " ")


CARD, LINES_TS = load_card_data()
LINES_AS_OF = _fmt_ts(LINES_TS) if LINES_TS else "Aug 23"
CODE2NAME = {v: k for k, v in NAME2CODE.items()}


def _american(p):
    """Fair (no-vig) American odds for win prob p."""
    if p >= 0.5:
        return f"–{round(100 * p / (1 - p))}"
    return f"+{round(100 * (1 - p) / p)}"


def _book_line(margin):
    """Model margin -> the half-point line a book would post (home persp.)."""
    x = round(margin * 2) / 2
    return x


def apply_card(g):
    """Overlay pipeline numbers onto a GAMES entry; returns the line-movement
    strip text (empty if the game isn't in card_data). Also sets g["move"],
    the compact first-seen → now form for the score bug."""
    c = CARD.get(g["cfbd"])
    if not c:
        return "", []
    ha, hb = NAME2CODE[c["home"]], NAME2CODE[c["away"]]
    disp = {"STAN": "Stanford", "BAY": "Baylor", "AUB": "Auburn",
            "CLEM": "Clemson", "LOU": "Louisville", "MISS": "Ole Miss",
            "WIS": "Wisconsin"}
    m = c["model_margin"]
    ph_raw = c["model_p_home"]
    # 1) the machine's line the way a book posts it (half points)
    line = _book_line(m)
    fav, num = (ha, line) if line >= 0 else (hb, -line)
    g["machine"] = (f"{disp.get(fav, fav)} –{num:g}" if num else "PK")
    g["raw_margin"] = f"{disp.get(ha, ha)} {m:+.1f}" if m >= 0 else f"{disp.get(hb, hb)} {-m:+.1f}"
    # 2) fair moneylines, both sides, from the model's win prob (no vig)
    g["fair"] = (f"{disp.get(ha, ha)} {_american(ph_raw)} · "
                 f"{disp.get(hb, hb)} {_american(1 - ph_raw)}")
    # 3) projected score = model margin laid over the MARKET total (the
    #    pipeline has no totals model; say so on the card)
    total = c.get("ou")
    if total is not None:
        hp, ap = (float(total) + m) / 2, (float(total) - m) / 2
        order = [(ha, hp), (hb, ap)] if hp >= ap else [(hb, ap), (ha, hp)]
        g["score"] = " – ".join(f"{disp.get(t, t)} {int(p + 0.5)}" for t, p in order)
        g["score_note"] = f"projected score · machine margin on the {float(total):g} market total"
    else:
        g["score"], g["score_note"] = "", ""
    books = c.get("books") or {}
    dk, bov = books.get("DraftKings", {}).get("spread"), books.get("Bovada", {}).get("spread")
    parts = [_dash(x) for x in (dk, bov) if x is not None]
    g["market"] = " / ".join(parts) if parts else c["mkt_spread"]
    mls = []
    for bname, short in (("DraftKings", "DK"), ("Bovada", "Bov")):
        b = books.get(bname) or {}
        if b.get("home_ml") is not None and b.get("away_ml") is not None:
            mls.append(f"{short} {disp.get(ha, ha)} {int(b['home_ml']):+d} / "
                       f"{disp.get(hb, hb)} {int(b['away_ml']):+d}")
    g["market_ml"] = " · ".join(mls)
    ph = round(ph_raw * 100)
    g["wp"] = (ha, ph, hb, 100 - ph)
    strip = ""
    first, now, clv = c.get("first_spread"), c.get("now_spread"), c.get("clv")
    # sides-flip guard: a first-seen with the opposite sign and an absurd CLV
    # is a known CFBD data error (e.g. Wisconsin-ND Jul 14) — fall back to the
    # book opener as the honest starting point and drop the corrupt CLV.
    flipped = (first is not None and now is not None and first * now < 0
               and clv is not None and abs(clv) > 20)
    if flipped:
        first, clv = c.get("opener"), None
    if first is not None:
        if flipped:
            strip = (f"opened {_dash(first)} → now {_dash(now)} "
                     "(first-seen excluded: sides-flip data error)")
        else:
            strip = (f"first-seen {_dash(first)} ({_fmt_ts(c['first_ts'])}) "
                     f"→ now {_dash(now)}")
            if c.get("opener") is not None and c["opener"] != first:
                strip += f" · opened {_dash(c['opener'])}"
            if clv is not None:
                strip += f" · CLV {clv + 0:+g}"
        g["move"] = f"{_dash(first)} → {_dash(now)}"
        if clv:
            g["move"] += f" · CLV {clv:+g}"
    return strip

_GAMES_EP1 = [  # Week 0 archive — the frozen Ep1 predictions (unrendered)
    dict(
        a="UNC", b="TCU", vs="vs", title="North Carolina vs TCU",
        cfbd=("North Carolina", "TCU"),
        where="Dublin, Ireland · Aviva Stadium",
        sub="Sat Aug 29 · noon ET · ESPN · Aer Lingus Classic · Belichick year two opens abroad",
        machine="TCU –1.5", market="–7.5 / –8", value="6–6.5 pts on UNC",
        wp=("TCU", 54, "UNC", 46),
        decides=[
            "UNC rush offense #127 vs TCU rush defense #28 — June may find nothing",
            "UNC's path: Edwards at TCU's #103 pass defense",
            "TCU pass offense #29 vs UNC pass defense #49 — best-on-best",
            "STRUCTURAL DISAGREEMENT: moneyline says TCU ~73%, machine says 54% — an 18-pt gap, not a spread quibble",
        ],
        honesty="Honesty check: UNC's −12.1 residual was 2025's biggest ACC market-overrating — the skepticism premium is earned. Line moved from −6.5 to −7.5 since Aug 18: the market is leaning further into TCU, not toward us.",
        lean="Lean: UNC ATS at −7.5 — research, not a play, while the win-prob gap is this wide · O/U no lean",
                keys_a=["Throw it — TCU's pass defense is the soft spot; the run game is closed", 'Win the hidden yards — the special-teams edge is real', "Bend-don't-break without Steve Belichick: keep the defense simple", 'Keep the circus outside: a clean, scripted first quarter'],
        keys_b=['Craig efficient, not heroic — distribute to a loaded receiver room', 'Make UNC one-dimensional with the run front, then rush the passer', 'Hold up in the back end with a new starting safety', 'Dwyer in the return game: the weapon inside a mediocre unit'],
        players_a=[("Billy Edwards Jr.", "QB — Maryland grad; 2,881 yds/15 TD in '24; '25 at Wisconsin lost to injury (7/16, 113 yds); Belichick's pick"),
                   ("Demon June", "RB — 2025: 464 yds at 5.5/carry, 3 total TD — Corey's key, but see the run-front matchup"),
                   ("Jordan Shipp", "WR — the X; 2025: 60 rec, 671 yds, 6 TD"),
                   ("Trech Kekahuna", "slot — 2025 at Wisconsin: 26 rec, 211 yds; 129 rush yds, 1 TD")],
        players_b=[("Jaden Craig", "QB — Harvard grad transfer; 2025: 61.6%, 2,829 yds, 25-7; first year as the guy post-Hoover"),
                   ("Jordan Dwyer", "WR — 2025: 54 rec, 730 yds, 7 TD (13.5 avg)"),
                   ("Ka'Morreun Pimpton", "TE — 2025: just 1 catch (a TD)"),
                   ("Markis Deal", "NT — anchors the #28 run front; 2025: 26 tkl, 2 TFL, 1 sack")],
    ),
    dict(
        a="NCSU", b="UVA", vs="at", title="NC State at Virginia",
        cfbd=("NC State", "Virginia"),
        where="Charlottesville · Scott Stadium",
        sub="Sat Aug 29 · 3:30 ET · revenge of the 35-31 end-zone INT · UVA: 11 wins, ACC CG loss in 2025",
        machine="UVA –6.7", market="–5.5 / –5.5", value="EDGE GONE — market converged",
        wp=("UVA", 65, "NCSU", 35),
        decides=[
            "The books converged: −3 / −6 on Aug 18 → −5.5 everywhere by Aug 22 — the edge we stated is gone",
            "UVA's #26 pass defense vs CJ Bailey's #36 pass offense decides it",
            "NC State explosiveness bottom-third: sustain or stall",
            "Special teams: UVA #56 vs NCSU #101 — ST cost State two 2025 games",
        ],
        honesty="UVA's +11.7 residual (11th-highest in FBS) says ESPN rates them well above what public inputs explain — and the roster lost 6 of its top 7 receivers and its senior pass rush. The machine's −6.7 is the stalest number on this card.",
        lean="No play: value was at −3, fair is −6, the market sits at −5.5 · machine still leans UVA",
                keys_a=['Bailey has to beat an elite pass defense with a rebuilt receiver room', 'Establish Duke Scott to stay out of obvious passing downs', 'Stop losing games on special teams', 'Make Pribula the story — pressure with the new edges'],
        keys_b=['Pribula ball security: no giveaways, no hero throws', 'Ride the veteran O-line and the run game', 'Replace the pass rush, not just the pass rushers', 'Let the pass defense win the rebuilt-receiver matchup'],
        players_a=[("CJ Bailey", "QB — yr 3; 2025: 3,105 yds, 25-9 at 68.8%; top-5 returning ACC QB"),
                   ("Duke Scott", "RB — 2025: 595 yds at 5.6/carry, 4 TD, plus 15 catches (Jayden on the stat sheet)"),
                   ("JoJo Trader", "WR — Miami transfer; 2025: 13 rec, 178 yds, 1 TD"),
                   ("Joseph Adedire", "DE — 2025: 6 tkl, 2 TFL, 1 sack")],
        players_b=[("Beau Pribula", "QB — Missouri transfer; 2025: 67.4%, 1,946 yds, 11-9 TD-INT (+6 rush TD); named starter Jul 15 at ACC Kickoff"),
                   ("Peyton Lewis", "RB — Tennessee transfer; 2025: 290 yds, 7 TD on 70 carries"),
                   ("Rico Flores Jr.", "WR — UCLA transfer; 2025: 26 rec, 274 yds"),
                   ("Fisher Camac", "DE — 2025: 44 tkl, 5.5 TFL, 3.5 sacks")],
    ),
    dict(
        a="JSU", b="NDSU", vs="at", title="Jacksonville State at North Dakota State",
        cfbd=("Jacksonville State", "North Dakota State"),
        where="Fargo · Fargodome",
        sub="Sat Aug 29 · 5:30 ET · CBSSN · NDSU's first FBS game · 2015 FCS title rematch · ON THE UPSET BOARD",
        machine="NDSU –2.7", market="–10 → –7", value="+3 CLV banked",
        wp=("NDSU", 57, "JSU", 43),
        decides=[
            "JSU identity: #23 rush offense, back-to-back ~3,500-yard rushing seasons",
            "JSU pass game #107 — stack the box, make Creel throw",
            "JSU special teams #123 — hidden-yardage leaks compound indoors",
            "NDSU has no ratings history — ESPN's −8.3 is a guess; max uncertainty",
        ],
        honesty="July 14 YEL alert: JSU +10, edge 7.3 — the market has since moved 3 points to our side (−7 at both books). MAX UNCERTAINTY: zero FBS ratings history, so the 57% deserves a range, not a point. And NDSU is FULLY ELIGIBLE — the NCAA repealed the transition ban Jun 24: MWC title game, bowls and the CFP are all open in year one.",
        lean="Model side: JSU at −7 or worse — already paid in line movement",
                keys_a=['Run into a front with no FBS tape — the identity must survive the RB change', "Stay ahead of the chains; don't make Creel throw 30 times", 'No hidden-yardage leaks indoors', 'Make Hayes beat you in his first start — disguise, third-and-long'],
        keys_b=['Protect Hayes early: scripted, short, on schedule', 'Lean on the run game and the O-line', "Tackle Creel in space — the rebuilt CB room's test", "Take the free points JSU's special teams give"],
        players_a=[("Caden Creel", "QB — 2025: 1,514 pass + 1,075 rush (5.9/carry), 16 total TD; dual-threat = Bison D's soft spot"),
                   ("Khristian Lando", "RB — the 'unleash' key; 2025: 201 yds on 51 carries behind Creel"),
                   ("Khurtiss Perry", "DT — 2025: 20 tkl, 4.5 TFL, 2 sacks"),
                   ("Jacob Cruz", "EDGE — 2025: 14 tkl, 1.5 TFL, 1 INT")],
        players_b=[("Nathan Hayes", "QB — FIRST career start; 2025 in relief: 25/44, 381 yds, 4-1, plus 88 rush yds; 'strongest arm since Wentz'"),
                   ("DJ Scott", "RB — 2025: 502 yds at 5.3/carry, 6 TD"),
                   ("Mekhi Collins", "WR — 2025: 6 catches, 159 yds, 2 TD — 26.5 per grab"),
                   ("Rebuilt CB room", "the one unit with no returning proof")],
    ),
    dict(
        a="HAW", b="STAN", vs="at", title="Hawai'i at Stanford",
        cfbd=("Hawai'i", "Stanford"),
        where="Palo Alto · Stanford Stadium",
        sub="Sat Aug 29 · 7:00 ET · ACC Network · Pritchard's first game · the machine rates Hawai'i the better team",
        machine="Stanford –1.6", market="–5.5 / –5.5", value="3.9 pts on Hawai'i",
        wp=("STAN", 54, "HAW", 46),
        decides=[
            "Neither side can run: Stanford rush offense #128 vs Hawai'i rush D #70; Hawai'i rush O #117 vs Stanford's #15 front",
            "So it's Alejado (#52 pass offense) at Stanford's #115 pass defense — the run-and-shoot's whole path",
            "Stanford QB Davis Warren: zero 2025 snaps at Michigan, career 7 TD–10 INT — unmodeled QB-tier risk",
            "Hawai'i's #8 special teams was Matsuzawa (27/29) + Barfield's return TD — the kicker is gone; unit rank overstates",
        ],
        honesty="Hawai'i's +5.8 residual: ESPN's 2025 rating ran ahead of its inputs, and the two sack leaders left. The line moved −3 → −5.5 toward Stanford since it opened — the market is not buying the machine's Hawai'i lean.",
        lean="Lean: Hawai'i +5.5 — a 4-pt gap, not a conviction; the opener at +3 was the worse number",
                keys_a=['Volume through the air at a bottom-20 pass defense', "Find the targets if Ashlock can't go", "Keep Alejado upright against Stanford's one good unit — the run front", "Three stops, not ten — and don't lose the kicking game with a new kicker"],
        keys_b=['Warren efficient, not expansive — first game back from the ACL', 'Make the run game exist', 'Tackle in space against four wide', 'Hold serve on special teams'],
        players_a=[("Micah Alejado", "QB — MWC Preseason POY; 2025: 66.3%, 3,106 yds, 24-9"),
                   ("Pofele Ashlock", "WR — 2025: 76 rec, 827 yds, 8 TD; the top target left (Jackson Harris, 963 yds)"),
                   ("Cam Barfield", "RB/KR — 2025: 371 rush yds, 4 TD; 28.7 per kick return with an 86-yd TD"),
                   ("Lesterlaisene Lagafuaina", "DL — 2025: 7.5 TFL, 3.5 sacks; the returning piece of a front that lost both sack leaders")],
        players_b=[("Davis Warren", "QB — Michigan transfer; 2025: no stat line (backup); career 5.4 yds/dropback, 7 TD–10 INT"),
                   ("Micah Ford", "RB — 2025: 643 yds at 4.4/carry, 4 TD, plus 11 catches"),
                   ("Nico Brown", "WR — Yale transfer; 2025: 71 rec, 1,085 yds, 11 TD in the Ivy"),
                   ("Matt Rose", "ILB — 2025: team-high 106 tkl, 8 TFL, 3 sacks")],
    ),
    dict(
        a="MEM", b="UNLV", vs="at", title="Memphis at UNLV",
        cfbd=("Memphis", "UNLV"),
        where="Las Vegas · Allegiant Stadium",
        sub="Sat Aug 29 · 10:00 ET · FOX · G5 heavyweight opener · Huff's 53-transfer debut · first film_study game",
        machine="UNLV –6.2", market="–5.5 / –5.5", value="market came to us: –3 → –5.5",
        wp=("UNLV", 64, "MEM", 36),
        decides=[
            "Memphis rush offense #5 vs UNLV rush defense #130 — but that #5 was Silverfield's roster; Huff flipped 53 players (111th returning production)",
            "UNLV rush offense #9 (Thomas, 7.0 a carry) vs Memphis rush D #20 — now led by transfers",
            "Arnold (#32 pass O context) at Memphis's #118 pass defense; Memphis's QB derby is a zero-FBS-snap tier",
            "Memphis's #16 special teams was Sutton Smith's 99-yd return TD — he's gone; UNLV #65",
        ],
        honesty="The machine's −6.2 is built on 2025 unit ranks for a Memphis roster that is 53 transfers new — both sides' numbers are stale in opposite directions. Residuals are flat (MEM −1.8, UNLV +0.7); the line already moved to our side.",
        lean="Model side UNLV — the edge was at −3 and it's gone; at −5.5 this is a watch and the first film_study run, not a play",
                keys_a=["Don't ask the (unannounced) quarterback to win it — run first", 'Run at the worst run defense in the league', 'Tapp has to get home on Arnold', 'Survive the hidden-yardage swing without Sutton Smith'],
        keys_b=["Feed Jai'Den Thomas", 'No turnovers from the quarterback, whoever it is', 'Make the new receivers real', "The run defense cannot be last year's unit against Huff-ball"],
        players_a=[("Air Noland / Marcus Stokes", "QB derby — Noland threw 3 passes at South Carolina in 2025; Stokes: 3,297 yds, 30 TD at D-II West Florida"),
                   ("Dallan Hayden", "RB — Colorado transfer; 2025: 326 yds at 4.7/carry, 1 TD"),
                   ("Tychaun Chapman", "WR — Southern Miss transfer; 2025: 24 rec, 444 yds, 3 TD"),
                   ("J'Mond Tapp", "DL — Southern Miss transfer; 2025: 69 tkl, 12 TFL, 7.5 sacks, 10 hurries")],
        players_b=[("Jackson Arnold", "QB — Auburn transfer; 2025: 63.3%, 1,309 yds, 6-2, +311 rush/8 TD in 8 starts (4-4); Orji still pushing"),
                   ("Jai'Den Thomas", "RB — MW POY candidate; 2025: 1,034 yds at 7.0/carry, 12 TD, +38 catches"),
                   ("Rebuilt WR room", "every 2025 target gone (Bradley 931, Omeire, Reynolds); Reddicks, Stellato, Walker via portal"),
                   ("Dee Crayton", "LB — Clemson transfer; 2025: 5 tkl as a reserve — the defense's 'patch' is unproven")],
    ),
]

GAMES = [
    dict(
        a="BAY", b="AUB", vs="vs", title="Baylor vs Auburn",
        cfbd=("Baylor", "Auburn"),
        where="Atlanta · Mercedes-Benz Stadium",
        sub="Sat Sep 5 · 3:30 ET · neutral site · Golesh's Auburn debut · Lagway's Baylor debut",
        machine="Auburn –5.5", market="–7 / –7", value="1.5 on Baylor · MONSTER UNDER under 59.5",
        wp=("AUB", 63, "BAY", 37),
        decides=[
            "Two rebooted programs, one neutral floor: Golesh's Auburn debut (USF's No. 2 total offense travels with him) vs Lagway's Baylor debut",
            "The starkest contrast on the card: Auburn returns 14% of its offense and rebuilt with 39 transfers — Baylor returns 60%, the most continuity here",
            "Lagway's health (shoulder and leg in 2025) is the Big 12 season's biggest single variable — Auburn's top-10 DL is the first test",
            "The 59.5 total is top-decile — five years of closing lines say that class goes UNDER 55% of the time, and two year-one offenses in a dome only helps",
        ],
        ctx_a=dict(coach="Aranda, year 6 — seat survived; new DC Klanderman (K-State)",
                   qb="NEW — DJ Lagway (Florida), biggest transfer get in program history",
                   roster="60% of offensive production back — most on this card · 31 portal adds"),
        ctx_b=dict(coach="NEW — Alex Golesh (USF) · Durkin RETAINED as DC",
                   qb="NEW and unsettled — open room after the Knight exit",
                   roster="14% back — most retooled roster on the card · 39 portal adds"),
        honesty="Both sidelines are running year-one installs the preseason numbers can't see — the machine's prior hasn't ingested a single 2026 snap. Auburn's edge is a talent prior; Baylor's counter is cohesion. Machine −5.5 vs market −7.5: a quibble, not a position.",
        keys_a=["Keep Lagway upright — his 2025 injury file is the season's whole risk profile; everything else is decoration",
                "Attack the rebuilt secondary: four transfer corners arrived — Auburn's front returned, the back end didn't",
                "Steal a possession with tempo — a halftime double-up covers the whole spread by itself",
                "Make the install-week offense chase: an early lead forces a unit that's never played together to hurry"],
        keys_b=["Lean on the defense that didn't change — 27 of the 39 transfers went to offense; Durkin's side kept its spine",
                "Tempo without turnovers: speed multiplies mistakes in week one of an install — take the give, live with punts",
                "Feed the tight ends where Golesh's 12-personnel system lives",
                "Win the line with top-10 DL talent — Lagway under duress is the whole ballgame"],
    ),
    dict(
        a="CLEM", b="LSU", vs="at", title="Clemson at LSU",
        cfbd=("Clemson", "LSU"),
        where="Baton Rouge · Tiger Stadium",
        sub="Sat Sep 5 · 7:30 ET · Kiffin's first game at LSU · Death Valley at night",
        machine="LSU –9", market="–10.5 / –10.5", value="1.5 to Clemson — quibble, not a play",
        wp=("LSU", 70, "CLEM", 30),
        decides=[
            "The biggest coaching debut of the season: Kiffin left a 13-win CFP team for this — first drive at night in Death Valley",
            "Dabo answers his worst season in 15 years with a throwback: Morris back as OC, first-year QB Vizzina — and Clemson opens UNRANKED vs No. 11",
            "Opposite repair jobs: LSU rebuilt with 44 transfers (21% back), Clemson added just 11 (41% back) — Saturday tells us which was right",
            "The market opened −11.5 and has walked to −10, toward the machine's −9 all week — and the loser burns a playoff mulligan",
        ],
        honesty="The machine's −9.1 is a preseason guess about a program running its second system and second staff inside one calendar year — and Clemson's offense is just as new. Install-week variance says treat every number here gently; 30–21 is a margin on a market total, not a script.",
        ctx_a=dict(coach="Dabo, year 18 — Chad Morris returns as OC",
                   qb="NEW — Vizzina, first-year starter post-Klubnik",
                   roster="41% back · just 11 portal adds — still portal-light by choice"),
        ctx_b=dict(coach="NEW — Lane Kiffin, off Ole Miss's 13-win CFP season",
                   qb="NEW — Sam Leavitt (Arizona State), the spring's biggest QB move",
                   roster="21% back · 44 portal adds — the biggest retool here"),
        keys_a=["Let Allen's defense set the terms — the only unit on this field that never stopped being itself",
                "Give Vizzina his pressure answers BEFORE the snap: quick game and screens, no improvising against Baker",
                "Win explosives — don't trade field goals with LSU's skill talent; the top-5 blue-chip roster has to cash as big plays",
                "Silence the crowd on third down — remove Death Valley's 2.5 and the machine has this at only 6.6"],
        keys_b=["Protect Leavitt — the portal-built OL on a short runway is the bet of the entire season",
                "Tempo early: force Allen's 4-2-5 to line up plain before the disguises load",
                "Let Baker hunt — SP+'s projected No. 2 defense against a first-career-start QB",
                "Finish in the red zone: settling for field goals keeps an unranked underdog alive all night"],
    ),
    dict(
        a="LOU", b="MISS", vs="vs", title="Louisville vs Ole Miss",
        cfbd=("Louisville", "Ole Miss"),
        where="Nashville · Nissan Stadium",
        sub="Sun Sep 6 · 7:30 ET · neutral site · Golding's debut as the head man · Chambliss back after the eligibility win",
        machine="Ole Miss –6.5", market="–7 / –6.5", value="machine = market — no gap",
        wp=("MISS", 65, "LOU", 35),
        decides=[
            "Golding's first game as the head man — can a promoted DC keep a 13-win CFP program's identity after Kiffin took its brain to Baton Rouge?",
            "Chambliss chose Oxford over following Kiffin — the rarest continuity win (50% back); Brohm counters with his third new QB1 in three years",
            "Ranked vs ranked, Sunday night in Nashville: #24 Louisville, #9 Ole Miss — neutral on paper, an Ole Miss crowd in practice",
            "SP+ projects Louisville favored in every game AFTER this one — win here and the huge ACC script is live; Ole Miss defends the CFP label",
        ],
        ctx_a=dict(coach="Brohm, year 4",
                   qb="NEW — Kienholz: third QB1 in three years, least-proven yet",
                   roster="35% back · 33 portal adds — the annual Brohm re-skin"),
        ctx_b=dict(coach="NEW — Pete Golding, promoted; kept the defensive spine",
                   qb="RETURNS — Chambliss, SEC Newcomer of the Year; won his eligibility appeal",
                   roster="50% back · 28 portal adds — the spine stayed"),
        honesty="The market opened −8.5 and walked to −7 — it has come to the machine's −6.5. That's agreement, not edge. And nothing in a preseason prior prices a first-year head coach's clock-and-timeout management on a Sunday neutral floor.",
        keys_a=["Ride the run game — the RB room is the roster's sneaky strength, and every carry keeps Chambliss on the sideline",
                "No coverage busts: top-20 havoc with bottom-40 explosives allowed is the Brohm tradeoff — one bust loses this game",
                "Pressure with the edges vs a rebuilt Ole Miss pass-pro — the one matchup the underdog clearly wins",
                "Keep Kienholz's menu short — Brohm's system has made a top-25 offense of every QB; let it work before asking for heroes"],
        keys_b=["Let Chambliss escape — statistically the nation's best at turning dead plays into first downs",
                "The interior line ends Louisville's run game early — force the least-proven Brohm QB to win it himself",
                "Tempo the Brohm defense into simple, static looks — no disguise time, no gambling looks",
                "Win field position all night — on a neutral floor against a new QB, the long field is a weapon"],
    ),
    dict(
        a="WIS", b="ND", vs="vs", title="Wisconsin vs Notre Dame",
        cfbd=("Wisconsin", "Notre Dame"),
        where="Green Bay · Lambeau Field",
        sub="Sun Sep 6 · 7:30 ET · neutral site · the 'Leave No Doubt' tour opens · Fickell's survival season",
        machine="ND –21", market="–20.5 / –20", value="no gap — watch, don't touch",
        wp=("ND", 87, "WIS", 13),
        decides=[
            "The revenge tour begins: 'Leave No Doubt' was born from the selection-show snub — Freeman rewatches the pain on purpose",
            "ND is the card's continuity outlier — Carr in year two, 51% of production back, just 7 portal adds — in a week full of reboots",
            "Fickell's survival season: a public vote of confidence, the Air Raid detour dead, and a G5 dual-threat QB making the B1G jump as a 20-point dog",
            "No. 4 ND plays for a top-seed résumé; Wisconsin's 3-3-5 — top-25 two straight years with zero offensive help — plays for a moral cover",
        ],
        ctx_a=dict(coach="Fickell, year 4 — survival season; the Air Raid detour is dead",
                   qb="NEW — Colton Joseph (Old Dominion), a real dual threat",
                   roster="31% back · 33 portal adds — retooled around a veteran OL"),
        ctx_b=dict(coach="Freeman — full staff and identity continuity",
                   qb="RETURNS — CJ Carr, year two: the top-10's likeliest QB leap",
                   roster="51% back · just 7 portal adds — the continuity program"),
        honesty="Housekeeping: our ledger's July first-seen here is a known CFBD sides-flip error — the honest movement read is the opener, −16.5 walked to −20.5, four points toward the Irish. Machine −21.1, market −20.5: agreement, and spreads this size sit in the favorite-longshot zone anyway. Watch, don't touch.",
        keys_a=["Shorten the game: long drives, clock runs, zero gifts — fewer possessions is the underdog's only math",
                "Find counters vs the tite front built to erase the run-first pivot — the league's most veteran OL has to earn it",
                "Joseph's legs on third-and-medium — the one chain-mover that travels up from the G5",
                "Force Carr to beat the 3-3-5 from the pocket — this defense is genuinely good enough to keep it inside 20"],
        keys_b=["Prove the post-Love duo/counter identity early — the gap-scheme run game sets up everything else",
                "Carr in rhythm: play-action on schedule, nothing forced — his year-two leap is the top 10's likeliest QB jump",
                "Erase explosives and make Wisconsin drive 12 plays — bend-don't-break with elite tackling wins by attrition",
                "'Leave No Doubt' means no scoreboard mercy — style points are seed points in a week-one résumé game"],
    ),
    dict(
        a="SMU", b="FSU", vs="at", title="SMU at Florida State",
        cfbd=("SMU", "Florida State"),
        where="Tallahassee · Doak Campbell Stadium",
        sub="Mon Sep 7 · 7:30 ET · Labor Day nightcap · ON THE UPSET BOARD — RED",
        machine="FSU –0.5", market="+3 / +3", value="RED: home dog outright",
        wp=("FSU", 52, "SMU", 48),
        decides=[
            "Monday night, alone on the calendar — the whole sport watches the week's last word",
            "No program swung harder in 36 months than FSU (CFP → 2–10 → rebound) — now it hosts SMU's post-CFP expectations and a No. 19 ranking",
            "SMU kept the roster AND Jennings (57% back) but lost BOTH coordinators to troikas; FSU kept Norvell but rents portal QBs a third time",
            "RED ALERT: machine takes the home dog outright — FSU 52% vs market 42%; strip home field and it actually favors SMU by 1.8",
        ],
        ctx_a=dict(coach="Lashlee, year 5 — both coordinators left; co-coordinator troikas",
                   qb="RETURNS — Kevin Jennings; the system's made 3 straight top-25 offenses",
                   roster="57% back · 15 portal adds — kept the roster, lost the callers"),
        ctx_b=dict(coach="Norvell, year 7 — Harris promoted to OC (Malzahn retired)",
                   qb="NEW — veteran portal room again",
                   roster="40% back · 22 portal adds — another portal reload"),
        honesty="Respect the backtest: RED alerts hit 49.7% ATS across 2023–25 — the flag is a research shortlist, not a pick. The entire machine case is the 2.5-point Monday-night home bump — neutralize the Doak and it likes SMU. And FSU's whiplash makes its prior the least trustworthy number on this card.",
        keys_a=["Prove the troika can call it — tempo and spacing from drive one, no committee hesitation on fourth downs",
                "Protect Jennings from FSU's top-10 portal front — line play is the one place SMU is a tier below",
                "Explosives, not field goals: White's 3-3-5 bends on purpose — long drives are the trap, chunk plays are the answer",
                "Win the takeaway ledger — the gambling back seven has to cash on the road, Monday night"],
        keys_b=["Run it with the portal front — make the game physical exactly where the trench classes graded top-10",
                "Quarterback run is the cheat code in Norvell's heavier-hand offense — the extra hat SMU's smaller front must answer",
                "No boom-bust: field position over hero ball — the whiplash program cannot beat itself Monday",
                "Start fast and make the Doak matter — the machine's whole case for FSU is worth 2.5 points of Monday night"],
    ),
]

# ---- Week 0 receipts (recap slide): frozen v12 predictions vs finals vs the
# last pre-kick ledger pull (Sat Aug 29, 10:16 ET). "off by" = distance of
# each side's line from the actual margin; mark M/K = machine/market closer.
RECAP = [
    ("UNC", "TCU", "North Carolina vs TCU",
     "we called TCU 25–23 · FINAL UNC 15–10",
     "our line TCU –1.5 · closing TCU –7.5",
     "machine off by 6.5 · market off by 12.5 — machine closer", "M"),
    ("NCSU", "UVA", "NC State at Virginia",
     "we called UVA 30–23 · FINAL UVA 34–8",
     "our line UVA –6.5 · closing UVA –4",
     "machine off by 19.5 · market off by 22 — machine closer", "M"),
    ("JSU", "NDSU", "Jacksonville State at NDSU",
     "we called NDSU 25–22 · FINAL NDSU 33–7",
     "our line NDSU –2.5 · closing NDSU –6.5",
     "machine off by 23.5 · market off by 19.5 — market closer", "K"),
    ("HAW", "STAN", "Hawai'i at Stanford",
     "we called Stanford 26–24 · FINAL Stanford 37–27",
     "our line Stanford –1.5 · closing Stanford –4",
     "machine off by 8.5 · market off by 6 — market closer", "K"),
    ("MEM", "UNLV", "Memphis at UNLV",
     "we called UNLV 32–26 · FINAL Memphis 27–21",
     "our line UNLV –6 · closing UNLV –4",
     "machine off by 12 · market off by 10 — market closer", "K"),
]

# ---- Superdog boards (segment: pick a dog to win outright; points = the
# spread). Computed live from card_data so a fresh pull refreshes them.
# AP Top 25 is week-1 2026 (CFBD /rankings) — update this dict each week.
AP_TOP25 = {"Ohio State": 1, "Oregon": 2, "Georgia": 3, "Notre Dame": 4,
            "Texas": 5, "Indiana": 6, "Miami": 7, "Texas A&M": 8,
            "Ole Miss": 9, "Oklahoma": 10, "LSU": 11, "Texas Tech": 12,
            "Alabama": 13, "USC": 14, "BYU": 14, "Michigan": 16,
            "Washington": 17, "Penn State": 18, "SMU": 19, "Tennessee": 20,
            "Utah": 21, "Iowa": 22, "Houston": 23, "Louisville": 24,
            "Missouri": 25}

_MONTHS = dict(Jan=1, Feb=2, Mar=3, Apr=4, May=5, Jun=6, Jul=7, Aug=8,
               Sep=9, Oct=10, Nov=11, Dec=12)


def superdog_boards():
    """(any-game rows, vs-top-25 rows), each sorted by EV; played games
    (before the lines_as_of date) are excluded."""
    import datetime as dt
    if not CARD:
        return [], []
    asof = dt.datetime.fromisoformat(LINES_TS.replace("Z", "+00:00")).date()
    rows = []
    for g in CARD.values():
        try:
            _, mon, day = g["date"].split(",")[0].split()
            gdate = dt.date(2026, _MONTHS[mon], int(day))
        except (KeyError, ValueError):
            continue
        if gdate < asof:
            continue
        b = g["books"].get("DraftKings") or g["books"].get("Bovada") or {}
        sp = b.get("spread")
        if sp is None or abs(sp) < 0.5:
            continue
        ph = g["model_p_home"]
        if sp < 0:
            dog, fav, p = g["away"], g["home"], 1 - ph
        else:
            dog, fav, p = g["home"], g["away"], ph
        at = "at" if dog == g["away"] else "vs"
        mkt_ph = g.get("mkt_p_home")
        mkt = (1 - mkt_ph if sp < 0 else mkt_ph) if mkt_ph is not None else None
        ml = b.get("away_ml") if dog == g["away"] else b.get("home_ml")
        rows.append(dict(dog=dog, fav=fav, at=at, pts=abs(sp), p=p, mkt=mkt,
                         ml=ml, ev=p * abs(sp), rank=AP_TOP25.get(fav)))
    rows.sort(key=lambda r: -r["ev"])
    return rows, [r for r in rows if r["rank"]]


SUPERDOG_ANY, SUPERDOG_T25 = superdog_boards()

# overlay pipeline numbers (card_data) before any slide is built
LEDGER = {g["title"]: apply_card(g) for g in GAMES}
print("card_data:", "loaded, lines as of " + LINES_AS_OF if CARD else
      "NOT FOUND — using hand-typed numbers")

# ---------------- title slide ----------------
s = blank(NAVY)
txt(s, 0.9, 0.85, 11.5, 0.45, "EPISODE 2 · WEEK 1 · SEP 5–7, 2026", 14, ORANGE,
    bold=True)
txt(s, 0.9, 1.2, 11.5, 1.1, "Five Games, One Card", 44, WHITE, bold=True)
y = 2.4
for g in GAMES:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 11.5, 0.8, NAVY2)
    logo_badge(s, 1.1, y + 0.1, 0.6, g["a"], plate=True)
    logo_badge(s, 1.85, y + 0.1, 0.6, g["b"], plate=True)
    txt(s, 2.7, y + 0.1, 6.5, 0.4, g["title"], 15, WHITE, bold=True)
    txt(s, 2.7, y + 0.46, 6.5, 0.3, g["where"], 10,
        RGBColor(0xCA, 0xDC, 0xFC))
    txt(s, 8.4, y + 0.1, 3.8, 0.4, g["machine"] + "   ·   " + g["market"],
        13.5, ORANGE, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, 8.4, y + 0.47, 3.8, 0.3, "machine · market", 9.5,
        RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.RIGHT)
    y += 0.88
txt(s, 0.9, 6.85, 11.5, 0.5,
    "Machine = ESPN 2026 preseason FPI + 2.5 HFA, empirical margin curve "
    f"(σ 17.9, 2021–25 fit) · lines as of {LINES_AS_OF} · model plays graded vs "
    "first-seen lines", 10.5,
    RGBColor(0x8F, 0xA5, 0xC4))

# ---------------- week 0 receipts ----------------
s = blank()
txt(s, 0.9, 0.5, 11.5, 0.55, "Week 0 — the receipts", 30, NAVY, bold=True)
txt(s, 0.9, 1.08, 11.5, 0.3,
    "Our call frozen at kickoff · closing line = last pre-kick pull "
    "(Sat 10:16 ET) · “off by” = miss vs the final margin",
    11, MUTE, italic=True)
y = 1.55
VERD = {"M": ORANGE, "K": RGBColor(0xB5, 0x12, 0x1B)}
for a, b, tit, callfin, lines_, miss, mark in RECAP:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 11.5, 0.86, ICE)
    logo_badge(s, 1.1, y + 0.13, 0.6, a)
    logo_badge(s, 1.8, y + 0.13, 0.6, b)
    txt(s, 2.6, y + 0.09, 4.8, 0.35, tit, 13.5, INK, bold=True)
    txt(s, 2.6, y + 0.46, 4.9, 0.3, callfin, 10.5, INK)
    txt(s, 7.35, y + 0.11, 4.8, 0.3, lines_, 9.5, MUTE)
    txt(s, 7.35, y + 0.44, 4.8, 0.3, miss, 10, VERD[mark], bold=True)
    y += 0.94
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y + 0.05, 11.5, 0.85, NAVY)
txt(s, 1.15, y + 0.17, 11.0, 0.65,
    "Total miss across five games: machine 70.0 points, closing market "
    "70.0 — a literal dead heat. Stated leans 1–2, the win coming where we "
    "disagreed with the market most (UNC, 18 points of win prob). Week 0 "
    "went UNDER in 4 of 5 · σ 17.9 — one week proves nothing either way.",
    10.5, WHITE)

# ---------------- per-game slides ----------------
for g in GAMES:
    # -- numbers slide --
    s = blank()
    logo_badge(s, 0.9, 0.42, 0.8, g["a"])
    txt(s, 1.82, 0.52, 0.5, 0.5, g["vs"], 14, MUTE, align=PP_ALIGN.CENTER)
    logo_badge(s, 2.35, 0.42, 0.8, g["b"])
    txt(s, 3.45, 0.38, 8.9, 0.55, g["title"], 27, NAVY, bold=True)
    txt(s, 3.45, 0.94, 8.9, 0.35, g["sub"], 11, MUTE, italic=True)

    # left: what decides it
    txt(s, 0.9, 1.75, 7.2, 0.4, "WHY IT MATTERS", 13, NAVY, bold=True)
    yy = 2.25
    for d in g["decides"]:
        shape(s, MSO_SHAPE.OVAL, 0.95, yy + 0.09, 0.14, 0.14, ORANGE)
        txt(s, 1.3, yy, 6.8, 0.8, d, 13.5)
        yy += 0.78
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, yy + 0.15, 7.2, 1.05, ICE)
    txt(s, 1.15, yy + 0.32, 6.7, 0.75, g["honesty"], 11.5, MUTE, italic=True)

    # right: navy score bug
    PALE = RGBColor(0xCA, 0xDC, 0xFC)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.5, 1.75, 3.9, 4.6, NAVY)
    txt(s, 8.8, 1.98, 3.3, 0.3, "THE NUMBER", 12, ORANGE, bold=True)
    # machine line as a book would post it + fair odds from our win prob
    txt(s, 8.8, 2.3, 3.3, 0.55, g["machine"], 28, WHITE, bold=True)
    txt(s, 8.8, 2.84, 3.3, 0.3, g.get("fair", ""), 12, WHITE, bold=True)
    txt(s, 8.8, 3.12, 3.3, 0.28,
        "machine line · fair odds, no vig · raw margin " + g.get("raw_margin", ""),
        8.5, PALE)
    # market
    txt(s, 8.8, 3.55, 3.3, 0.45, g["market"], 20, WHITE, bold=True)
    txt(s, 8.8, 3.98, 3.3, 0.4, "market (DK / Bovada) · " + g.get("market_ml", ""),
        8.5, PALE)
    # score prediction (replaced "the gap" per Lucas's in-Slides edit 8/31;
    # the authored value/gap strings stay in GAMES as data)
    score = g.get("score", "")
    txt(s, 8.8, 4.45, 3.3, 0.4, score, 15 if len(score) <= 22 else 13,
        ORANGE, bold=True)
    txt(s, 8.8, 4.82, 3.3, 0.25, "score prediction", 8.5, PALE)
    wa, pa, wb, pb = g["wp"]
    wp_bar(s, 8.8, 5.2, 3.3, 0.4, wa, pa, TEAMS[wa]["color"],
           wb, pb, TEAMS[wb]["color"])
    txt(s, 8.8, 5.66, 3.3, 0.25, "win probability (machine)", 8.5, PALE)
    if g.get("move"):
        txt(s, 8.8, 5.95, 3.3, 0.3, "LINE MOVE  " + g["move"], 9.5,
            ORANGE, bold=True)
    strip = LEDGER[g["title"]]
    if strip:
        txt(s, 0.9, 6.82, 11.5, 0.3, "Line movement: " + strip, 10.5, MUTE)

    # -- team slides (Corey's format: one full slide per team; the score
    #    predictions live only on the closing card, truly LAST). Context band
    #    = the significance frame: coach status, QB situation, roster
    #    continuity (CFBD returning offensive PPA + portal-add counts). --
    for key, keys, ctx in ((g["a"], g["keys_a"], g["ctx_a"]),
                           (g["b"], g["keys_b"], g["ctx_b"])):
        s = blank()
        logo_badge(s, 0.9, 0.55, 1.15, key)
        txt(s, 2.35, 0.66, 9.9, 0.62, CODE2NAME[key], 32, NAVY, bold=True)
        txt(s, 2.35, 1.34, 9.9, 0.35, "What it takes to win · " + g["title"],
            13, MUTE, italic=True)
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 1.95, 11.5, 1.05, ICE)
        for i, (label, val) in enumerate((("COACH", ctx["coach"]),
                                          ("QB", ctx["qb"]),
                                          ("ROSTER", ctx["roster"]))):
            x = 1.15 + i * 3.85
            txt(s, x, 2.08, 3.55, 0.25, label, 9.5, ORANGE, bold=True)
            txt(s, x, 2.34, 3.55, 0.6, val, 10.5, INK)
        yy = 3.4
        for k in keys:
            shape(s, MSO_SHAPE.OVAL, 1.0, yy + 0.12, 0.15, 0.15, ORANGE)
            txt(s, 1.45, yy, 10.5, 0.8, k, 15.5, INK)
            yy += 0.92
        txt(s, 0.9, 7.13, 11.5, 0.3, "EP 2 · WEEK 1 · " + g["title"], 9, MUTE)

# ---------------- closing card: predictions + superdogs ----------------
s = blank(NAVY)
PALE = RGBColor(0xCA, 0xDC, 0xFC)
txt(s, 0.9, 0.55, 11.5, 0.45, "EPISODE 2 · THE CARD", 14, ORANGE, bold=True)
txt(s, 0.9, 0.95, 11.5, 0.8, "Our predictions", 36, WHITE, bold=True)
y = 1.95
for g in GAMES:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 11.5, 0.72, NAVY2)
    logo_badge(s, 1.1, y + 0.1, 0.52, g["a"], plate=True)
    logo_badge(s, 1.75, y + 0.1, 0.52, g["b"], plate=True)
    txt(s, 2.5, y + 0.19, 5.4, 0.4, g["title"], 14, WHITE, bold=True)
    txt(s, 7.0, y + 0.09, 5.2, 0.42, g.get("score", ""), 15.5, ORANGE,
        bold=True, align=PP_ALIGN.RIGHT)
    txt(s, 7.0, y + 0.47, 5.2, 0.25, "market " + g["market"], 8.5,
        RGBColor(0x8F, 0xA5, 0xC4), align=PP_ALIGN.RIGHT)
    y += 0.8
# superdog band
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y + 0.08, 11.5, 1.02, ORANGE)
for i, (label, board) in enumerate(
        [("SUPERDOG", [r for r in SUPERDOG_ANY if not r["rank"]]),
         ("GIANT KILLER", SUPERDOG_T25)]):
    if not board:
        continue
    r = board[0]
    fav = (f"#{r['rank']} " if r["rank"] else "") + r["fav"]
    ml = f" · ML {int(r['ml']):+d}" if r.get("ml") is not None else ""
    txt(s, 1.2, y + 0.17 + i * 0.44, 3.0, 0.35, "★ " + label, 13, NAVY,
        bold=True)
    txt(s, 3.6, y + 0.17 + i * 0.44, 8.6, 0.35,
        f"{r['dog']} +{r['pts']:g} {r['at']} {fav}{ml}", 13.5, WHITE,
        bold=True)
txt(s, 0.9, 7.18, 11.5, 0.3,
    "projected scores = machine margin on the market total · superdogs = "
    "dog to win outright, points = the spread · graded vs first-seen lines "
    "· research, not picks", 9, RGBColor(0x8F, 0xA5, 0xC4), italic=True)

out = r"C:\Users\lucas\Fun Projects\Sports Data Analysis\ncaa-fbs-model\decks\2026_Week1_Episode2.pptx"
prs.save(out)
print("wrote", out, f"- {len(prs.slides)} slides")
